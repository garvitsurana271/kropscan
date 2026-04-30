"""Generate 4-slide pitch deck as PPTX"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

BG_DARK = RGBColor(10, 10, 10)
GREEN = RGBColor(61, 111, 57)
GREEN_LIGHT = RGBColor(34, 197, 94)
WHITE = RGBColor(255, 255, 255)
WHITE_DIM = RGBColor(150, 150, 150)
RED = RGBColor(239, 68, 68)
ORANGE = RGBColor(249, 115, 22)


def set_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_text(slide, left, top, width, height, text, font_size=18, color=WHITE, bold=False, align=PP_ALIGN.LEFT, font_name='Calibri'):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = align
    return txBox


# ═══════════════════════════════════════
# SLIDE 1: LOGO
# ═══════════════════════════════════════
slide1 = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
set_bg(slide1, RGBColor(15, 26, 15))

add_text(slide1, 0, 0.3, 13.333, 0.5, 'Case 7 \u2014 AI Crop Health Assistant',
         font_size=12, color=WHITE_DIM, align=PP_ALIGN.RIGHT)

add_text(slide1, 0, 2.2, 13.333, 1.2, 'KropScan',
         font_size=72, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

add_text(slide1, 0, 3.5, 13.333, 0.6, 'Offline-First AI Crop Disease Detection',
         font_size=24, color=WHITE_DIM, align=PP_ALIGN.CENTER)

add_text(slide1, 0, 4.8, 13.333, 0.5, 'Garvit Surana & Zaviaa  \u2022  Assam, India',
         font_size=16, color=WHITE_DIM, align=PP_ALIGN.CENTER)

# Feature pills
pills = ['156 Diseases', '20 Crops', '95.91% Accuracy', '11 Languages', 'Works Offline']
pill_text = '     \u2022     '.join(pills)
add_text(slide1, 0, 5.8, 13.333, 0.4, pill_text,
         font_size=13, color=GREEN_LIGHT, align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════
# SLIDE 2: PROBLEM
# ═══════════════════════════════════════
slide2 = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide2, BG_DARK)

add_text(slide2, 0.8, 0.5, 12, 1.5, '\u20b990,000 Cr',
         font_size=96, color=RED, bold=True)

add_text(slide2, 0.8, 2.0, 10, 0.8, 'Lost to crop disease every year in India',
         font_size=28, color=WHITE_DIM)

# Stats row
stats = [
    ('5 days', 'Average wait for diagnosis'),
    ('1 : 1,100', 'Agro officer to farmer ratio'),
    ('40%', 'Rural India without reliable internet'),
]
for i, (num, label) in enumerate(stats):
    x = 0.8 + i * 4.2
    add_text(slide2, x, 3.5, 3.8, 0.8, num,
             font_size=48, color=WHITE, bold=True)
    add_text(slide2, x, 4.3, 3.8, 0.6, label,
             font_size=13, color=WHITE_DIM)

add_text(slide2, 0.8, 5.8, 12, 0.8,
         '\u201cNot because it can\u2019t be treated \u2014 because farmers can\u2019t get a diagnosis in time.\u201d',
         font_size=18, color=RGBColor(100, 100, 100))

# ═══════════════════════════════════════
# SLIDE 3: NUMBERS
# ═══════════════════════════════════════
slide3 = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide3, RGBColor(10, 15, 10))

add_text(slide3, 0.8, 0.4, 5, 0.4, 'KROPSCAN BY THE NUMBERS',
         font_size=13, color=GREEN_LIGHT, bold=True)

# Grid of stats
grid = [
    ('~2 Lakh', 'Training Images'),
    ('156', 'Disease Classes'),
    ('20', 'Crops Covered'),
    ('95.91%', 'Validation Accuracy'),
    ('11', 'Indian Languages'),
    ('3 sec', 'Diagnosis Time (Offline)'),
]
for i, (num, label) in enumerate(grid):
    col = i % 3
    row = i // 3
    x = 0.8 + col * 4.2
    y = 1.2 + row * 2.0
    add_text(slide3, x, y, 3.5, 0.7, num,
             font_size=40, color=GREEN_LIGHT, bold=True)
    add_text(slide3, x, y + 0.65, 3.5, 0.4, label,
             font_size=12, color=WHITE_DIM)

# Revenue breakdown
add_text(slide3, 0.8, 5.2, 5, 0.3, 'REVENUE MODEL',
         font_size=11, color=GREEN_LIGHT, bold=True)

rev = [
    ('38%', 'Brand Recommendations'),
    ('36%', 'API & Data Licensing'),
    ('24%', 'Pro Subscriptions (\u20b949/mo)'),
    ('2%', 'E-Commerce Referrals'),
]
for i, (pct, name) in enumerate(rev):
    x = 0.8 + i * 3.2
    add_text(slide3, x, 5.6, 2.8, 0.5, pct,
             font_size=28, color=GREEN_LIGHT, bold=True)
    add_text(slide3, x, 6.1, 2.8, 0.4, name,
             font_size=10, color=WHITE_DIM)

# ═══════════════════════════════════════
# SLIDE 4: MVP VIDEO
# ═══════════════════════════════════════
slide4 = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide4, BG_DARK)

add_text(slide4, 0, 0.5, 13.333, 0.4, 'MVP DEMO',
         font_size=13, color=GREEN_LIGHT, bold=True, align=PP_ALIGN.CENTER)

# Video placeholder
add_text(slide4, 0, 2.5, 13.333, 1, '\u25b6  Insert screen recording here',
         font_size=32, color=WHITE_DIM, align=PP_ALIGN.CENTER)

add_text(slide4, 0, 3.6, 13.333, 0.5, 'Airplane Mode \u2192 Scan \u2192 Result \u2192 Hindi \u2192 Market Prices \u2192 Chatbot',
         font_size=14, color=RGBColor(80, 80, 80), align=PP_ALIGN.CENTER)

# Feature pills at bottom
features = '\U0001f512 Offline AI (WASM)     \U0001f30d 11 Languages     \U0001f4ca Live Mandi Prices     \U0001f916 AI Chatbot     \U0001f48a 1-Tap Treatment     \U0001f4f1 PWA + Android'
add_text(slide4, 0, 5.8, 13.333, 0.4, features,
         font_size=12, color=GREEN_LIGHT, align=PP_ALIGN.CENTER)

prs.save('docs/KropScan_Pitch_Deck.pptx')
print('PPTX pitch deck saved')
