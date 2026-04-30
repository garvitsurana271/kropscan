"""
Fill pitch_template.pptx with KropScan-NER content for HackDays 4.0 submission.

Strategy: walk every text frame in the template, detect placeholder patterns
(text with "[Describe...", "INSERT...", "ADD YOUR..."), and replace the whole
frame's content with the corresponding KropScan slide content. Preserves the
template's font/color via the first run's formatting.

Usage:
    python scripts/fill_pitch_deck.py

Output:
    docs/KropScan_NER_Pitch.pptx
"""
import sys
import io
import os
from copy import deepcopy
from pptx import Presentation
from pptx.util import Pt
from pptx.enum.text import PP_ALIGN

sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8')

ROOT = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
TEMPLATE = os.path.join(ROOT, 'pitch_template.pptx')
OUTPUT = os.path.join(ROOT, 'docs', 'KropScan_NER_Pitch.pptx')


# ---------------------------------------------------------------------------
# Slide content. Each entry maps a "trigger phrase" found in the template
# slide to the replacement text. The trigger phrase identifies which text
# frame holds the placeholder; the replacement is the new full-frame content.
# ---------------------------------------------------------------------------
SLIDE_CONTENT = [
    # Slide 1 — Title (special-cased; replaces multiple frames in field_replacements)
    {
        'mode': 'fields',
        'field_replacements': {
            'Problem Statement Title:': 'Problem Statement Title:\nKropScan NER — Hybrid Edge-AI for Marginal NE Farmers',
            'Theme:': 'Theme:\n#2 Agriculture / 2.5 AI-Based Crop Recommendation & Advisory',
            'Team Name:': 'Team Name:\n[Your Team Name]',
            'Team Members': 'Team Members:\nGarvit Surana (Full-stack + AI), Mehul Jain (Frontend Designer)',
            'college Name:': 'College Name:\n[Your College Name]',
        }
    },
    # Slide 2 — Problem
    {
        'mode': 'replace',
        'triggers': ['Core Pain Point', 'Target Audience', 'Current Flaws', 'ADD YOUR PROBLEM'],
        'content': (
            "CORE PAIN POINT\n"
            "Marginal NE farmers lose ₹15K-50K per acre to crop disease outbreaks "
            "they can't diagnose in time. Every existing app needs internet — and "
            "none cover NER specialty crops like areca nut, large cardamom, king "
            "chilli, or Khasi mandarin.\n\n"
            "TARGET AUDIENCE\n"
            "• 11M+ marginal NER farmers (1-2 acres, intermittent connectivity)\n"
            "• ASHA workers and KVK officers needing district-level surveillance\n\n"
            "CURRENT FLAWS\n"
            "• Plantix and gov apps are cloud-only — break in NER villages\n"
            "• Zero existing apps cover NER specialty crops\n"
            "• No advisory layer for jhum (shifting cultivation) erosion crisis"
        )
    },
    # Slide 3 — Solution
    {
        'mode': 'replace',
        'triggers': ['Core Concept', 'Value Proposition', 'EXPLAIN YOUR SOLUTION', 'high-level summary'],
        'content': (
            "CORE CONCEPT\n"
            "Three-tier hybrid AI for offline-first NE farming:\n\n"
            "1. LOCAL — 81MB ONNX model in WebAssembly. ~95% on 40 common crops. "
            "Cached in IndexedDB. Fully offline.\n"
            "2. CLOUD — Gemma 4 Vision fallback for NER specialty crops (areca "
            "nut, large cardamom, ginger, turmeric, black rice, king chilli, "
            "kiwi, Khasi mandarin, tea, bamboo).\n"
            "3. GRACEFUL FALLBACK — never crashes; clear retry UX.\n\n"
            "Plus: jhum rotation advisory + ASHA disease surveillance + 15 "
            "languages incl. 4 NE tribal langs (Bodo / Mizo / Khasi / Manipuri).\n\n"
            "VALUE PROPOSITION\n"
            "• Connectivity-resilient: works where every other ag app fails\n"
            "• Region-specific: only app covering NER crops + tribal languages\n"
            "• Public-health tier: outbreak surveillance for state ag departments"
        )
    },
    # Slide 4 — Tech Stack
    {
        'mode': 'replace',
        'triggers': ['LIST YOUR TECH STACK', 'Insert UI Frameworks', 'Insert Server', 'Insert DB'],
        'content': (
            "FRONTEND\n"
            "• React 19 + TypeScript + Vite 6\n"
            "• Tailwind CSS 3\n"
            "• Capacitor 8 (Android wrapper, installable PWA)\n"
            "• Recharts (ASHA heatmap)\n\n"
            "BACKEND / AI\n"
            "• ONNX Runtime Web (in-browser WASM inference)\n"
            "• EfficientNetV2-S (40-class plant disease classifier)\n"
            "• Google Gemma 4 31B (cloud vision fallback)\n"
            "• Gemini 2.5 Flash (UI translation engine)\n"
            "• Optional Flask service for server-side inference\n\n"
            "DATABASE & CLOUD\n"
            "• Firebase Phone OTP Auth + Firestore real-time DB\n"
            "• IndexedDB (offline scan history, ONNX model cache)\n"
            "• Workbox Service Worker (offline-first PWA)"
        )
    },
    # Slide 5 — System Architecture
    {
        'mode': 'replace',
        'triggers': ['Data Flow & Integrations', 'INSERT SYSTEM ARCHITECTURE', 'how data flows'],
        'content': (
            "DATA FLOW\n"
            "Camera → Image Quality Gate → Local ONNX (WASM)\n"
            "   ↓ (if confidence < 0.55 OR blight)\n"
            "Gemma 4 Vision (cloud) → Diagnosis → Treatment Plan in ₹/acre\n"
            "   ↓\n"
            "Stored to IndexedDB (offline) + Firestore (sync)\n\n"
            "KEY ALGORITHMS\n"
            "• Hybrid escalation with confidence + blight-specific override\n"
            "• Severity-weighted aggregation: score = Σ severity × log10(farmers+1)\n"
            "• Slope/soil erosion model calibrated to ICAR-NEH research\n\n"
            "PERFORMANCE & SECURITY\n"
            "• Bundle 1.5MB JS gzipped; 81MB ONNX cached forever in IndexedDB\n"
            "• First scan <3s cold, <1s warm\n"
            "• Phone OTP auth (no email leak vector); ASHA reports anonymized\n"
            "• ~₹0.05 per user-month on Gemma free tier"
        )
    },
    # Slide 6 — Prototype Demo
    {
        'mode': 'replace',
        'triggers': ['DESCRIBE PROTOTYPE', 'DEMO HERE', 'Key Demo Features', 'Highlight'],
        'content': (
            "WHAT JUDGES WILL SEE (~3 minutes live)\n\n"
            "1. SCAN AN ARECA NUT LEAF — local ONNX returns low confidence → "
            "app silently escalates to Gemma → \"Areca Nut Yellow Leaf Disease "
            "(Phytoplasma)\" + ₹2,200/acre treatment plan in INR. Pull WiFi "
            "mid-scan to prove offline capability for the 40 common crops.\n\n"
            "2. JHUM ROTATION ADVISORY — pick crop / years cleared / slope / "
            "soil → instant \"Critical erosion risk 75/100. Rotate Rice → Soybean "
            "→ Maize → Finger Millet → 8-yr fallow.\" AI agronomist insights load.\n\n"
            "3. SWITCH UI TO BODO — Settings → Language → Bodo. Sidebar, "
            "dashboard, scan flow all flip to Bodo Devanagari. Cached forever.\n\n"
            "WOW FACTORS\n"
            "• Fully offline crop diagnosis in browser via WebAssembly\n"
            "• NER specialty crops without retraining (hybrid prompt-eng fallback)\n"
            "• Tribal language UI auto-translated and cached"
        )
    },
    # Slide 7 — Roadmap
    {
        'mode': 'replace',
        'triggers': ['Immediate Next Steps', 'Core Feature Expansion', 'Beta Launch',
                     'Long-term Vision', 'OUTLINE YOUR FUTURE'],
        'content': (
            "PHASE 1 — IMMEDIATE (next 2 weeks)\n"
            "• Field-validate NER crop diagnoses with KVK officers\n"
            "• Source 500+ real NER disease photos for retraining set\n"
            "• Bake tribal language translations at build time\n\n"
            "PHASE 2 — CORE EXPANSION (1-2 months)\n"
            "• Retrain ONNX on NER-specific dataset\n"
            "• Bluetooth/WiFi-Direct mesh sync for offline ASHA reports\n"
            "• Native Android APK distribution via KVK partnerships\n\n"
            "PHASE 3 — BETA LAUNCH (3-6 months)\n"
            "• 1,000-farmer pilot in 3 NE districts\n"
            "• Integration with Govt of India ENAM mandi prices\n"
            "• Voice-input chatbot for low-literacy farmers (Whisper API)\n\n"
            "PHASE 4 — LONG-TERM (6-18 months)\n"
            "• B2G licensing of ASHA dashboard to state ag departments\n"
            "• Expansion to other connectivity-poor regions\n"
            "• Crop insurance partnership integration"
        )
    },
    # Slide 8 — Thank You
    {
        'mode': 'fields',
        'field_replacements': {
            'Add GitHub Link': 'github.com/garvitsurana271/kropscan',
            'Add Team Lead Email': 'dev@409.ai',
        }
    },
]


def clear_text_frame(text_frame):
    """Clear all paragraphs from a text frame (keep one empty placeholder paragraph)."""
    # Keep first paragraph but clear its runs
    if not text_frame.paragraphs:
        return
    first = text_frame.paragraphs[0]
    # Clear all runs in first paragraph
    for r in list(first.runs):
        r.text = ''
    # Remove all subsequent paragraphs by manipulating XML
    p_element = first._p
    parent = p_element.getparent()
    siblings = list(parent)
    keep_idx = siblings.index(p_element)
    for sib in siblings[keep_idx + 1:]:
        if sib.tag == p_element.tag:
            parent.remove(sib)


def set_text_frame_content(text_frame, content):
    """Replace text frame content with `content` (multi-line string).
    Preserves font from the first run if available."""
    # Capture font from first non-empty run before clearing
    sample_run_font = None
    for para in text_frame.paragraphs:
        for run in para.runs:
            if run.text.strip():
                sample_run_font = {
                    'name': run.font.name,
                    'size': run.font.size,
                    'bold': run.font.bold,
                    'color_rgb': run.font.color.rgb if run.font.color and run.font.color.type is not None else None,
                }
                break
        if sample_run_font:
            break

    clear_text_frame(text_frame)

    lines = content.split('\n')
    if not lines:
        return

    # First line goes into the existing first paragraph
    first_para = text_frame.paragraphs[0]
    if not first_para.runs:
        run = first_para.add_run()
    else:
        run = first_para.runs[0]
    run.text = lines[0]
    if sample_run_font:
        if sample_run_font['name']:
            run.font.name = sample_run_font['name']
        if sample_run_font['size']:
            try:
                run.font.size = sample_run_font['size']
            except Exception:
                pass

    # Subsequent lines as new paragraphs
    for line in lines[1:]:
        new_para = text_frame.add_paragraph()
        new_run = new_para.add_run()
        new_run.text = line
        if sample_run_font:
            if sample_run_font['name']:
                new_run.font.name = sample_run_font['name']
            if sample_run_font['size']:
                try:
                    new_run.font.size = sample_run_font['size']
                except Exception:
                    pass


def text_frame_text(text_frame):
    return '\n'.join(
        ''.join(r.text for r in para.runs) for para in text_frame.paragraphs
    )


def fill_template():
    print(f'Loading template: {TEMPLATE}')
    prs = Presentation(TEMPLATE)
    print(f'Slides found: {len(prs.slides)}')
    print()

    for idx, slide in enumerate(prs.slides):
        if idx >= len(SLIDE_CONTENT):
            print(f'Slide {idx+1}: no content defined, skipping')
            continue
        spec = SLIDE_CONTENT[idx]
        slide_changes = 0

        # Collect all text frames from this slide
        text_frames = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                text_frames.append((shape, shape.text_frame))

        if spec['mode'] == 'fields':
            # Per-frame field-style replacements (used for title slide and thank-you)
            for shape, tf in text_frames:
                full = text_frame_text(tf)
                for trigger, replacement in spec['field_replacements'].items():
                    if trigger in full:
                        set_text_frame_content(tf, replacement)
                        slide_changes += 1
                        break
        else:  # 'replace' mode — find the largest body frame matching any trigger and replace it
            triggers = spec['triggers']
            content = spec['content']
            # Find the LARGEST text frame (by char count) that matches any trigger.
            # This is heuristically the body content area, not a header label.
            candidates = []
            for shape, tf in text_frames:
                full = text_frame_text(tf)
                if any(t.lower() in full.lower() for t in triggers):
                    candidates.append((len(full), shape, tf))
            if candidates:
                candidates.sort(key=lambda c: c[0], reverse=True)
                _, target_shape, target_tf = candidates[0]
                set_text_frame_content(target_tf, content)
                slide_changes += 1

                # AGGRESSIVELY clear EVERY other text frame on this slide that
                # looks like a leftover instruction. We check against a broad
                # set of patterns since placeholder phrasing varies.
                placeholder_patterns = [
                    '[', 'Insert', 'Describe', 'Explain ', 'Highlight', 'Mention',
                    'Detail any', 'Detail], ', 'INSERT', 'DESCRIBE', 'EXPLAIN',
                    'ADD YOUR', 'LIST YOUR', 'DEMO HERE', 'OUTLINE YOUR',
                    'PROTOTYPE', 'unique benefit', 'high-level summary',
                    'sentence]', 'matters]', 'missing]', 'innovative?',
                    'under the hood?', 'database and back]',
                    'integrated]', 'optimizations]',
                    'Vue.js', 'FastAPI', 'PostgreSQL, AWS',
                    'Bug Fixes]', 'Acquisition]', 'Scaling]',
                    'Polish]', 'Add GitHub Link', 'Add Team Lead Email',
                ]
                target_id = id(target_tf)
                for shape, tf in text_frames:
                    if id(tf) == target_id:
                        continue
                    full = text_frame_text(tf)
                    if not full.strip():
                        continue
                    # Only clear if text matches a known placeholder pattern
                    if any(kw in full for kw in placeholder_patterns):
                        clear_text_frame(tf)
                        slide_changes += 1

        print(f'Slide {idx+1}: {slide_changes} frame(s) updated')

    os.makedirs(os.path.dirname(OUTPUT), exist_ok=True)
    prs.save(OUTPUT)
    print(f'\n[OK] Wrote {OUTPUT}')


if __name__ == '__main__':
    fill_template()
